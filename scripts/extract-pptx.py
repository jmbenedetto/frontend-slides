#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx",
# ]
# ///
"""Extract text, images, and notes from a PowerPoint file."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Extract slide text, images, and notes from a .pptx file.",
    )
    parser.add_argument("input_pptx", type=Path, help="Path to the source .pptx file.")
    parser.add_argument(
        "output_dir",
        nargs="?",
        type=Path,
        default=Path("."),
        help="Directory for extracted assets and extracted-slides.json.",
    )
    return parser.parse_args()


def extract_pptx(file_path: Path, output_dir: Path) -> list[dict[str, object]]:
    presentation = Presentation(file_path)
    slides_data: list[dict[str, object]] = []

    assets_dir = output_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_title = slide.shapes.title.text if slide.shapes.title is not None else ""
        slide_data: dict[str, object] = {
            "number": slide_index,
            "title": slide_title,
            "content": [],
            "images": [],
            "notes": "",
        }

        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                text_content = shape.text.strip()
                if text_content:
                    slide_data["content"].append(
                        {"type": "text", "content": text_content}
                    )

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_name = (
                    f"slide{slide_index}_img{len(slide_data['images']) + 1}.{image.ext}"
                )
                image_path = assets_dir / image_name
                image_path.write_bytes(image.blob)

                slide_data["images"].append(
                    {
                        "path": f"assets/{image_name}",
                        "width": shape.width,
                        "height": shape.height,
                    }
                )

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            slide_data["notes"] = notes_frame.text.strip()

        slides_data.append(slide_data)

    return slides_data


def main() -> int:
    args = parse_args()
    input_pptx = args.input_pptx.expanduser().resolve()
    output_dir = args.output_dir.expanduser().resolve()

    if not input_pptx.is_file():
        raise SystemExit(f"Input file not found: {input_pptx}")

    output_dir.mkdir(parents=True, exist_ok=True)
    slides = extract_pptx(input_pptx, output_dir)

    output_path = output_dir / "extracted-slides.json"
    with output_path.open("w", encoding="utf-8") as handle:
        json.dump(slides, handle, indent=2, ensure_ascii=False)
        handle.write("\n")

    print(f"Extracted {len(slides)} slides to {output_path}")
    for slide in slides:
        image_count = len(slide["images"])
        title = slide["title"] or "(no title)"
        print(f"  Slide {slide['number']}: {title} — {image_count} image(s)")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
